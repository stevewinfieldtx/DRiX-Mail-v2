import csv, io, json, re
from pathlib import Path
from fastapi import FastAPI, Depends, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from sqlalchemy import select
from sqlalchemy.orm import Session, joinedload
from pypdf import PdfReader
from pptx import Presentation
from openpyxl import load_workbook
from .config import settings
from .database import Base, engine, get_db
from .models import Client, Campaign, Prospect, Evidence, Email, ClientStatus, ProspectStatus
from .schemas import ClientCreate, ClientOut, CampaignCreate, CampaignOut, ProspectCreate, ProspectOut, EmailOut, ProfileUpdate, ReplyIn
from .engine import build_seller_profile, research, generate_next, mark_sent, classify_reply, apply_reply

app=FastAPI(title="B2B Narrative Platform",version="1.0.0")
app.add_middleware(CORSMiddleware,allow_origins=settings.cors_origins.split(","),allow_credentials=True,allow_methods=["*"],allow_headers=["*"])

@app.on_event("startup")
def startup(): Base.metadata.create_all(engine)

@app.get("/health")
def health(): return {"status":"ok"}

@app.get("/api/clients",response_model=list[ClientOut])
def clients(db:Session=Depends(get_db)): return db.scalars(select(Client).order_by(Client.name)).all()

@app.post("/api/clients",response_model=ClientOut)
def create_client(data:ClientCreate,db:Session=Depends(get_db)):
    c=Client(**data.model_dump()); db.add(c); db.commit(); db.refresh(c); return c

@app.delete("/api/clients/{client_id}")
def delete_client(client_id:str,db:Session=Depends(get_db)):
    c=db.get(Client,client_id)
    if not c: raise HTTPException(404,"Client not found")
    if c.status!=ClientStatus.draft: raise HTTPException(409,"Only draft clients can be deleted")
    if c.campaigns: raise HTTPException(409,"Draft client has campaigns and cannot be deleted")
    db.delete(c); db.commit()
    return {"deleted":client_id}

def extract_upload(upload:UploadFile)->str:
    raw=upload.file.read(); ext=Path(upload.filename or "").suffix.lower()
    if ext==".pdf": return "\n".join((p.extract_text() or "") for p in PdfReader(io.BytesIO(raw)).pages)
    if ext in {".pptx",".ppt"}:
        deck=Presentation(io.BytesIO(raw)); return "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape,"text_frame"))
    return raw.decode("utf-8",errors="ignore")

@app.post("/api/clients/{client_id}/sources",response_model=ClientOut)
def add_sources(client_id:str,files:list[UploadFile]=File(...),db:Session=Depends(get_db)):
    c=db.get(Client,client_id)
    if not c: raise HTTPException(404,"Client not found")
    c.source_texts=[*c.source_texts,*[{"filename":f.filename,"text":extract_upload(f)[:100000]} for f in files]]; db.commit(); return c

@app.post("/api/clients/{client_id}/generate-profile",response_model=ClientOut)
def profile(client_id:str,db:Session=Depends(get_db)):
    c=db.get(Client,client_id)
    if not c: raise HTTPException(404,"Client not found")
    c.profile=build_seller_profile(c); c.status=ClientStatus.draft; db.commit(); return c

@app.put("/api/clients/{client_id}/approve",response_model=ClientOut)
def approve(client_id:str,data:ProfileUpdate,db:Session=Depends(get_db)):
    c=db.get(Client,client_id)
    if not c: raise HTTPException(404,"Client not found")
    c.profile=data.profile; c.status=ClientStatus.approved; db.commit(); return c

@app.get("/api/campaigns",response_model=list[CampaignOut])
def campaigns(db:Session=Depends(get_db)): return db.scalars(select(Campaign)).all()

@app.post("/api/campaigns",response_model=CampaignOut)
def create_campaign(data:CampaignCreate,db:Session=Depends(get_db)):
    c=db.get(Client,data.client_id)
    if not c or c.status!=ClientStatus.approved: raise HTTPException(409,"Seller profile must be approved first")
    x=Campaign(**data.model_dump()); db.add(x); db.commit(); db.refresh(x); return x

@app.get("/api/prospects",response_model=list[ProspectOut])
def prospects(status:str|None=None,db:Session=Depends(get_db)):
    q=select(Prospect).order_by(Prospect.created_at.desc())
    if status: q=q.where(Prospect.status==status)
    return db.scalars(q).all()

@app.post("/api/prospects",response_model=ProspectOut)
def create_prospect(data:ProspectCreate,db:Session=Depends(get_db)):
    if not db.get(Campaign,data.campaign_id): raise HTTPException(404,"Campaign not found")
    p=Prospect(**data.model_dump()); db.add(p); db.commit(); db.refresh(p); return p

@app.post("/api/prospects/{prospect_id}/research",response_model=ProspectOut)
def research_prospect(prospect_id:str,db:Session=Depends(get_db)):
    p=db.scalar(select(Prospect).options(joinedload(Prospect.campaign).joinedload(Campaign.client)).where(Prospect.id==prospect_id))
    if not p: raise HTTPException(404,"Prospect not found")
    return research(db,p)

@app.get("/api/prospects/{prospect_id}/evidence")
def ledger(prospect_id:str,db:Session=Depends(get_db)): return db.scalars(select(Evidence).where(Evidence.prospect_id==prospect_id)).all()

@app.post("/api/prospects/{prospect_id}/generate-next",response_model=EmailOut)
def next_email(prospect_id:str,db:Session=Depends(get_db)):
    p=db.scalar(select(Prospect).options(joinedload(Prospect.emails),joinedload(Prospect.campaign).joinedload(Campaign.client)).where(Prospect.id==prospect_id))
    if not p: raise HTTPException(404,"Prospect not found")
    if p.status==ProspectStatus.new: p=research(db,p)
    try:return generate_next(db,p)
    except ValueError as e: raise HTTPException(409,str(e))

@app.post("/api/emails/{email_id}/sent",response_model=ProspectOut)
def sent(email_id:str,db:Session=Depends(get_db)):
    e=db.scalar(select(Email).options(joinedload(Email.prospect).joinedload(Prospect.campaign)).where(Email.id==email_id))
    if not e: raise HTTPException(404,"Email not found")
    return mark_sent(db,e)

@app.post("/api/prospects/{prospect_id}/reply",response_model=ProspectOut)
def reply(prospect_id:str,data:ReplyIn,db:Session=Depends(get_db)):
    p=db.scalar(select(Prospect).options(joinedload(Prospect.emails)).where(Prospect.id==prospect_id))
    if not p: raise HTTPException(404,"Prospect not found")
    kind=data.kind or classify_reply(data.text); return apply_reply(db,p,kind)

IMPORT_FIELDS=["company_url","company_name","contact_name","contact_title","contact_email","linkedin_url","additional_urls","notes","external_campaign_id"]
SYNONYMS={"company_url":["companyurl","website","url","domain","companywebsite","webaddress"],"company_name":["companyname","company","organization","organisation","account","businessname"],"contact_name":["contactname","fullname","name","decisionmaker","prospectname"],"contact_title":["contacttitle","title","jobtitle","role","position"],"contact_email":["contactemail","email","emailaddress","workemail"],"linkedin_url":["linkedinurl","linkedin","linkedinprofile"],"additional_urls":["additionalurls","extraurls","otherurls","sources"],"notes":["notes","context","comments","description"],"external_campaign_id":["externalcampaignid","externalid","campaignid"]}
def _normalized(value): return re.sub(r"[^a-z0-9]","",str(value or "").lower())
def _xlsx_columns(row):
    columns=[]; seen={}
    for index,value in enumerate(row):
        header=str(value or "").strip()
        if not header:continue
        seen[header]=seen.get(header,0)+1
        label=header if seen[header]==1 else f"{header} ({seen[header]})"
        columns.append((index,label))
    return columns
def _tabular(upload:UploadFile):
    raw=upload.file.read(); ext=Path(upload.filename or "").suffix.lower()
    if ext==".xlsx":
        sheet=load_workbook(io.BytesIO(raw),read_only=True,data_only=True).active
        values=list(sheet.iter_rows(values_only=True))
        if not values:return [],[]
        columns=_xlsx_columns(values[0]); headers=[label for _,label in columns]
        rows=[{label:(row[index] if index<len(row) and row[index] is not None else "") for index,label in columns} for row in values[1:] if any(index<len(row) and row[index] is not None and str(row[index]).strip() for index,_ in columns)]
        return headers,rows
    if ext not in {".csv",".txt"}: raise HTTPException(422,"Upload a .csv or .xlsx file")
    reader=csv.DictReader(io.StringIO(raw.decode("utf-8-sig")))
    return list(reader.fieldnames or []),list(reader)
def _suggest(headers):
    normalized={_normalized(h):h for h in headers}; result={}
    for field,aliases in SYNONYMS.items():
        match=next((normalized[a] for a in aliases if a in normalized),None)
        if match:result[field]=match
    return result

@app.post("/api/batch/inspect")
def batch_inspect(file:UploadFile=File(...)):
    headers,rows=_tabular(file)
    return {"filename":file.filename,"headers":headers,"suggested_mapping":_suggest(headers),"sample":rows[:3],"row_count":len(rows)}

@app.post("/api/batch/import/{campaign_id}")
def batch_import(campaign_id:str,file:UploadFile=File(...),mapping:str=Form("{}"),db:Session=Depends(get_db)):
    if not db.get(Campaign,campaign_id): raise HTTPException(404,"Campaign not found")
    headers,rows=_tabular(file)
    try:column_map=json.loads(mapping)
    except json.JSONDecodeError:raise HTTPException(422,"Invalid column mapping")
    column_map={field:source for field,source in column_map.items() if field in IMPORT_FIELDS and source in headers}
    if not column_map:column_map={field:field for field in IMPORT_FIELDS if field in headers}
    created=[]; errors=[]
    used=set(column_map.values())
    for n,row in enumerate(rows,start=2):
        try:
            value=lambda field:str(row.get(column_map.get(field,""),"") or "").strip()
            if not value("company_url") and not value("company_name"):raise ValueError("Map either Company URL or Company Name")
            known={k:value(k) for k in ["company_url","company_name","contact_name","contact_title","contact_email","linkedin_url","notes","external_campaign_id"]}
            known["extra_urls"]=[u.strip() for u in re.split(r"[|,\n]",value("additional_urls")) if u.strip()]
            known["custom_fields"]={k:v for k,v in row.items() if k not in used and v not in (None,"")}
            p=Prospect(campaign_id=campaign_id,**known); db.add(p); db.flush(); created.append(p.id)
        except Exception as exc: errors.append({"row":n,"error":str(exc)})
    db.commit(); return {"created":len(created),"prospect_ids":created,"errors":errors,"mapping":column_map}

EXPORT_FIELDS=["prospect_id","company_name","company_url","contact_name","contact_title","contact_email","campaign_name","campaign_stage","email_number","subject","body","cta","confidence_score","narrative_type","primary_business_outcome","proof_used","research_summary","generated_at","status","send_after","stop_if_replied","thread_subject","sequence_id","external_campaign_id"]
@app.get("/api/batch/export/{campaign_id}")
def batch_export(campaign_id:str,generate:bool=True,db:Session=Depends(get_db)):
    ps=db.scalars(select(Prospect).options(joinedload(Prospect.emails),joinedload(Prospect.campaign).joinedload(Campaign.client)).where(Prospect.campaign_id==campaign_id)).unique().all(); out=io.StringIO(); w=csv.DictWriter(out,fieldnames=EXPORT_FIELDS); w.writeheader()
    for p in ps:
        if p.status in {ProspectStatus.closed,ProspectStatus.suppressed,ProspectStatus.meeting,ProspectStatus.replied}: continue
        e=next((x for x in p.emails if x.status.value=="ready"),None)
        if not e and generate:
            if p.status==ProspectStatus.new: research(db,p)
            if p.status!=ProspectStatus.waiting or not p.due_at or p.due_at.isoformat() <= __import__('datetime').datetime.now(__import__('datetime').timezone.utc).isoformat():
                try:e=generate_next(db,p)
                except ValueError:continue
        if not e: continue
        w.writerow({"prospect_id":p.id,"company_name":p.company_name,"company_url":p.company_url,"contact_name":p.contact_name,"contact_title":p.contact_title,"contact_email":p.contact_email,"campaign_name":p.campaign.name,"campaign_stage":p.stage,"email_number":e.number,"subject":e.subject,"body":e.body,"cta":e.cta,"confidence_score":p.confidence,"narrative_type":e.narrative_type,"primary_business_outcome":e.primary_business_outcome,"proof_used":e.proof_used,"research_summary":e.research_summary,"generated_at":e.generated_at.isoformat(),"status":e.status.value,"send_after":p.due_at.isoformat() if p.due_at else "","stop_if_replied":True,"thread_subject":e.thread_subject,"sequence_id":e.sequence_id,"external_campaign_id":p.external_campaign_id})
    out.seek(0); return StreamingResponse(iter([out.getvalue()]),media_type="text/csv",headers={"Content-Disposition":"attachment; filename=next-emails.csv"})

@app.get("/api/dashboard")
def dashboard(db:Session=Depends(get_db)):
    ps=db.scalars(select(Prospect).options(joinedload(Prospect.campaign).joinedload(Campaign.client)).order_by(Prospect.due_at.nullsfirst())).all()
    return [{"prospect_id":p.id,"client":p.campaign.client.name,"campaign":p.campaign.name,"prospect":p.company_name or p.company_url,"stage":p.stage,"last_action":max([e.sent_at or e.generated_at for e in p.emails],default=p.created_at),"reply_status":p.reply_status,"confidence":p.confidence,"due_date":p.due_at,"next":p.next_action,"status":p.status.value} for p in ps]
